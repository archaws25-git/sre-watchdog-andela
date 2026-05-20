"""Generate the SRE Watchdog MVP presentation as a PowerPoint file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
        tf.paragraphs[i].level = 0
        tf.paragraphs[i].font.size = Pt(18)


def add_screenshot_slide(prs, title, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    # Placeholder box for screenshot
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.3), Inches(8), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
    shape.line.color.rgb = RGBColor(0x36, 0xA2, 0xEB)
    # Text inside placeholder
    tf2 = shape.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf2.paragraphs[0].text = "[INSERT SCREENSHOT HERE]"
    tf2.paragraphs[0].font.size = Pt(20)
    tf2.paragraphs[0].font.color.rgb = RGBColor(0x36, 0xA2, 0xEB)
    # Caption
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(8), Inches(0.5))
    tf3 = txBox2.text_frame
    tf3.text = caption
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.italic = True
    tf3.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "SRE Watchdog",
        "AI-Powered Intelligent Observability & Event Watchdog\n"
        "MVP Presentation — May 2026"
    )

    # Slide 2: Problem Statement
    add_content_slide(prs, "The Problem", [
        "SRE teams drown in log noise — thousands of entries per minute",
        "Manual threshold monitoring misses subtle degradation patterns",
        "Alert fatigue from duplicate/noisy notifications",
        "No intelligent context in alerts — just raw numbers",
        "Need: AI-powered anomaly detection with actionable summaries",
    ])

    # Slide 3: Solution Overview
    add_content_slide(prs, "The Solution: SRE Watchdog", [
        "Python FastAPI application with AI-powered log analysis",
        "Two-gate detection: statistical pre-filter + AWS Bedrock AI",
        "Automatic webhook alerts with severity classification",
        "Real-time dashboard with Chart.js visualizations",
        "Full audit trail — every detection, analysis, and alert recorded",
        "12-factor app methodology — all config via environment variables",
    ])

    # Slide 4: Architecture
    add_content_slide(prs, "Architecture Overview", [
        "Gate 1: APScheduler tick → per-service error rate computation",
        "Gate 2: FastAPI BackgroundTask → Bedrock Converse API (Claude Sonnet 4.6)",
        "Alert Service: Severity mapping + webhook dispatch with retry",
        "Cooldown: Prevents alert fatigue (15-min suppression window)",
        "Dashboard: Jinja2 SSR + Chart.js + auto-refresh every 60s",
        "Storage: SQLite WAL mode for concurrent read/write",
    ])

    # Slide 5: Tech Stack
    add_content_slide(prs, "Tech Stack", [
        "Language: Python 3.11+",
        "Framework: FastAPI + Uvicorn",
        "AI: AWS Bedrock (Claude Sonnet 4.6 via Converse API)",
        "Database: SQLite with WAL mode (SQLAlchemy ORM)",
        "Scheduling: APScheduler (BackgroundScheduler)",
        "Dashboard: Jinja2 + Chart.js",
        "Testing: pytest + hypothesis + freezegun + respx",
    ])

    # Slide 6: Development Timeline
    add_content_slide(prs, "Development Timeline", [
        "Phase 1: Requirements & Design (spec-driven development)",
        "Phase 2: Project scaffold, config, database foundation",
        "Phase 3: Core services (ingestion, Bedrock client, alert service)",
        "Phase 4: Anomaly detector (two-gate pipeline) + scheduler",
        "Phase 5: API routers (8 endpoints) + app wiring",
        "Phase 6: Dashboard, synthetic log generator, documentation",
        "Phase 7: Test infrastructure, verification, git setup",
    ])

    # Slide 7: Key Features
    add_content_slide(prs, "Key Features Delivered", [
        "10,000 synthetic logs across 5 services (24-hour simulation)",
        "3 seeded anomaly windows: sharp spike, sustained degradation, cascade",
        "AI-generated anomaly summaries with severity scoring (0.0–1.0)",
        "Webhook alerts with 4-band severity (LOW/MEDIUM/HIGH/CRITICAL)",
        "Cooldown suppression to prevent alert fatigue",
        "On-demand analysis via POST /analyze + job polling",
        "11 API endpoints + HTML dashboard",
    ])

    # Slide 8: Dashboard Screenshot - Overview
    add_screenshot_slide(
        prs,
        "Dashboard: Metrics & Error Rate Chart",
        "Metrics bar showing logs ingested, anomalies, alerts, and failures. "
        "Chart.js line chart with error rate per service over 24 hours."
    )

    # Slide 9: Dashboard Screenshot - Anomalies
    add_screenshot_slide(
        prs,
        "Dashboard: Anomaly Detection Results",
        "Recent anomalies table showing service, time window, AI score, "
        "lifecycle status, severity badge, and AI-generated summary."
    )

    # Slide 10: Dashboard Screenshot - Alerts
    add_screenshot_slide(
        prs,
        "Dashboard: Alert Dispatch Log",
        "Recent alerts table showing timestamp, service, severity, "
        "anomaly ID, and dispatch status (sent/suppressed/failed)."
    )

    # Slide 11: Dashboard Screenshot - Run Analysis
    add_screenshot_slide(
        prs,
        "Dashboard: On-Demand Analysis",
        "Run Analysis button triggering Bedrock AI analysis across all services. "
        "Shows loading state and refreshes results on completion."
    )

    # Slide 12: Detection Pipeline
    add_content_slide(prs, "Detection Pipeline in Action", [
        "1. APScheduler tick fires every 60 seconds",
        "2. Gate 1: Compute error rate per service (sliding 5-min window)",
        "3. If error_rate > 10%: create AnomalyWindow (pending_analysis)",
        "4. Gate 2: BackgroundTask invokes Bedrock with log context",
        "5. Bedrock returns anomaly_score (0.0–1.0) + AI summary",
        "6. If score ≥ 0.5 and no cooldown: dispatch webhook alert",
        "7. Full lifecycle persisted for audit trail",
    ])

    # Slide 13: AI Integration
    add_content_slide(prs, "AWS Bedrock AI Integration", [
        "Model: Claude Sonnet 4.6 (us.anthropic.claude-sonnet-4-6)",
        "Prompt: Service context + error rate + log messages (capped at 50)",
        "Response: JSON with anomaly_score (0.0–1.0) + plain-text summary",
        "Retry: Exponential backoff (1s, 2s, 4s) for transient errors",
        "Health caching: /health reports last-known Bedrock status",
        "Cost control: Gate 1 pre-filter reduces unnecessary API calls",
    ])

    # Slide 14: Challenges & Solutions
    add_content_slide(prs, "Challenges & Solutions", [
        "Model ID changes: Iterated through 4 model IDs to find active one",
        "Response parsing: Model wraps JSON in markdown — added fence stripping",
        "Cooldown noise: Suppressed records cluttered dashboard — by design for audit",
        "Concurrent writes: SQLite WAL mode enables BackgroundTask parallelism",
        "Session credentials: Temporary tokens expire — documented in workflow",
    ])

    # Slide 15: Code Quality
    add_content_slide(prs, "Code Quality & Standards", [
        "48 source files, 7,322 lines of code",
        "flake8 linting: zero errors (max-line-length=120)",
        "Google-style docstrings on all public functions",
        "Module-level docstrings in every Python file",
        "Structured JSON logging throughout",
        "Type hints on all function signatures",
        "Comprehensive .env.example with all 14 configuration variables",
    ])

    # Slide 16: What's Next
    add_content_slide(prs, "Production Roadmap", [
        "Authentication: OAuth2/OIDC for dashboard access",
        "Deployment: Docker + AWS App Runner or ECS/Fargate",
        "Semantic analysis: S3 Vectors for log embeddings",
        "Cursor-based pagination for high-volume log stores",
        "Test coverage: Complete the optional test suite (target 80%+)",
        "Monitoring: CloudWatch integration for self-observability",
    ])

    # Slide 17: Demo
    add_content_slide(prs, "Live Demo", [
        "1. Start server: uvicorn app.main:app --reload",
        "2. Generate logs: python generate_logs.py",
        "3. Open dashboard: http://localhost:8000/dashboard",
        "4. Click 'Run Analysis' — watch Bedrock AI in action",
        "5. Observe: anomaly scores, AI summaries, alert dispatch",
        "6. Check /health, /metrics, /anomalies endpoints",
    ])

    # Slide 18: Thank You
    add_title_slide(
        prs,
        "Thank You",
        "SRE Watchdog — AI-Powered Observability\n"
        "Questions?"
    )

    output_path = "SRE_Watchdog_MVP_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
